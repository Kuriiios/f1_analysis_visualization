import seaborn as sns
from matplotlib import pyplot as plt
from datetime import timedelta

from instagrapi import Client
import subprocess
from pdf2image import convert_from_path

import fastf1
import fastf1.plotting
from fastf1.ergast import Ergast

import os
import csv

import pandas as pd
import numpy as np
import ast

from pathlib import Path
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from datetime import datetime

ergast = Ergast()

from utils.race.lap_info import *
from utils.race.plots import *

parent_file = Path(__file__).resolve().parent.parent

plots_dir = os.path.join(parent_file, 'plots')
if plots_dir not in sys.path:
    sys.path.append(plots_dir)
    print(f"Added {plots_dir} to sys.path")

cache_folder = parent_file / 'cache_folder'
cache_folder.mkdir(exist_ok=True)

fastf1.Cache.enable_cache(cache_folder)

fastf1.plotting.setup_mpl(mpl_timedelta_support=True, misc_mpl_mods=False,
                          color_scheme='fastf1')

year = datetime.today().year
race_number = fastf1.get_events_remaining()['RoundNumber'].values[0]
race_session = 'R'
race_session_name = 'Race'

session = fastf1.get_session(year, race_number, race_session)
session.load()

event_name = session.event.EventName   
figures_folder = parent_file / 'figures' / f'{race_number}_{event_name}_{session.event.year}/'
report_folder = parent_file / 'reports' / f"{race_number}_{event_name}_{session.event.year}/"

figures_folder.mkdir(parents=True, exist_ok=True)
report_folder.mkdir(parents=True, exist_ok=True)

pit = ergast.get_pit_stops(season = year, round = race_number )
circuit_info = session.get_circuit_info()
teams = fastf1.plotting.list_team_names(session)

def read_credentials(file = parent_file / "data/raw/login.txt"):
    creds = {}
    with open(file) as f:
        for line in f:
            key, value = line.strip().split("=")
            creds[key] = value
    return creds

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))
    
def Hex_RGB(ip):
    return tuple(int(ip[i+1:i+3],16) for i in (0, 2, 4))

def to_str(var):
    return str(list(np.reshape(np.asarray(var), (1, np.size(var)))[0]))[1:-1]

nat_as_integer = np.datetime64('NAT').view('i8')

def isnat(your_datetime):
    dtype_string = str(your_datetime.dtype)
    if 'datetime64' in dtype_string or 'timedelta64' in dtype_string:
        return your_datetime.view('i8') == nat_as_integer
    return False  # it can't be a NaT if it's not a dateime

def correct_fastest_lap(team, session):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    
    driver_1_laps = session.laps.pick_drivers(team_drivers[0]).pick_laps(range(0, (int(max(session.laps['LapNumber'])) + 1))).reset_index()
    driver_2_laps = session.laps.pick_drivers(team_drivers[1]).pick_laps(range(0, (int(max(session.laps['LapNumber'])) + 1))).reset_index()
    if not driver_1_laps.empty:
        laptime_counter_driver_1 = 0
        for lap in driver_1_laps['LapTime']:
            try:
                if 'NaT' in str(lap):
                    driver_1_laps.loc[laptime_counter_driver_1, 'LapTime'] = driver_1_laps.loc[laptime_counter_driver_1 +1, 'LapStartTime'] - driver_1_laps.loc[laptime_counter_driver_1, 'LapStartTime']
                laptime_counter_driver_1 +=1
            except KeyError:
                break

    laptime_counter_driver_2 = 0
    for lap in driver_2_laps['LapTime']:
        try:
            if 'NaT' in str(lap):
                driver_2_laps.loc[laptime_counter_driver_2, 'LapTime'] = driver_2_laps.loc[laptime_counter_driver_2 +1, 'LapStartTime'] - driver_2_laps.loc[laptime_counter_driver_2, 'LapStartTime']
            laptime_counter_driver_2 +=1
        except KeyError:
            break
    return driver_1_laps, driver_2_laps

for idx,team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    team_color = fastf1.plotting.get_team_color(team, session=session)
    df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
    team_color_2 = df_color.iat[idx,0]
    
    show_pace_comp(team, session, team_color, team_color_2, figures_folder)
    show_laptime_scatterplot(team, team_drivers, session, team_color, team_color_2, figures_folder)
    show_laptime_comp(team, session, team_color, team_color_2, figures_folder)
    show_tyre_strategy(team, team_drivers, session, figures_folder)
    

laptime_scatterplot = {}
laptime_comp = {}
tyre_strategy = {}
pace_comp0 = {}
pace_comp1 = {}

csv_file_path = f'{race_number}_{race_session}_race_info.csv'
csv_file_path_fastest_driver = f'{race_number}_{race_session}_race_fastest_driver.csv'

lap_info_per_team = [['EventName', 'driver_1_name', 'driver_1_position', 'driver_1_gap', 'total_time_driver_1', 'avg_laptime_driver_1','iqr_driver_1', 'number_pit_driver_1', 'total_duration_pit_driver_1', 'driver_2_name', 'driver_2_position', 'driver_2_gap', 'total_time_driver_2', 'avg_laptime_driver_2','iqr_driver_2', 'number_pit_driver_2', 'total_duration_pit_driver_2', 'fastest_driver_1', 'fastest_driver_2', 'safety_car_lap', 'fastest_lap']]
race_info = []
fastest_driver_per_lap_dict = {}
for idx,team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    drivers_info = get_drivers_info(session, team, team_drivers, race_session, year, race_number)
    fastest_driver_per_lap_per_team = get_faster_driver_per_lap(session, team, team_drivers)
    fastest_driver_per_lap_dict.update({f'{team}':fastest_driver_per_lap_per_team})
    team_info = get_lap_repartition( fastest_driver_per_lap_per_team)
    race_info = create_csv_race_info(lap_info_per_team, drivers_info, team_info, event_name, race_session)
os.chdir(parent_file / 'data/processed/')
with open(csv_file_path, mode='w', newline='') as file:
    writer = csv.writer(file)
    writer.writerows(race_info)

keyword = f'{race_number}_{race_session}_race_info'
for fname in os.listdir(parent_file / 'data/processed/'):
    if keyword in fname:
        driver_data = fname
os.chdir(parent_file / 'data/processed/')
arr = pd.read_csv(driver_data)

counter = 0
prs = Presentation()

for idx, team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    team_color = fastf1.plotting.get_team_color(team, session=session)
    df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
    team_color_2 = df_color.iat[idx,0]
    team_color = Hex_RGB(team_color)
    team_color_2 = Hex_RGB(team_color_2)
    
    try:
        race_name = arr.iloc[counter]['EventName']
        team_logo = parent_file / f'data/external/team_logos/{team}.png'
        laptime_comp = figures_folder / f'{race_session_name}_{team}_laptime_comp.png'
        laptime_scatterplot = figures_folder / f'{race_session_name}_{team}_laptime_scatterplot.png'
        tyre_strategy = figures_folder / f'{race_session_name}_{team}_tyre_strategy.png'
        driver_1_pace = figures_folder / f'{race_session_name}_{team}_driver_1_pace.png'
        driver_2_pace = figures_folder / f'{race_session_name}_{team}_driver_2_pace.png'
        
        name_driver_1 = arr.iloc[counter]['driver_1_name']
        driver_1_position = arr.iloc[counter]['driver_1_position']
        driver_1_gap = arr.iloc[counter]['driver_1_gap']
        total_time_driver_1 = arr.iloc[counter]['total_time_driver_1']
        avg_laptime_driver_1 = arr.iloc[counter]['avg_laptime_driver_1']
        iqr_driver_1 = arr.iloc[counter]['iqr_driver_1']
        number_pit_driver_1 = arr.iloc[counter]['number_pit_driver_1']
        total_duration_pit_driver_1 = arr.iloc[counter]['total_duration_pit_driver_1']
        
        driver_2_name = arr.iloc[counter]['driver_2_name']
        driver_2_position = arr.iloc[counter]['driver_2_position']
        driver_2_gap = arr.iloc[counter]['driver_2_gap']
        total_time_driver_2 = arr.iloc[counter]['total_time_driver_2']
        avg_laptime_driver_2 = arr.iloc[counter]['avg_laptime_driver_2']
        iqr_driver_2 = arr.iloc[counter]['iqr_driver_2']
        number_pit_driver_2 = arr.iloc[counter]['number_pit_driver_2']
        total_duration_pit_driver_2 = arr.iloc[counter]['total_duration_pit_driver_2']
        
        fastest_driver_1 = arr.iloc[counter]['fastest_driver_1']
        fastest_driver_2 = arr.iloc[counter]['fastest_driver_2']
        safety_car_lap = arr.iloc[counter]['safety_car_lap']
        fastest_lap = arr.iloc[counter]['fastest_lap']
        
        prs.slide_width = Pt(1080)
        prs.slide_height = Pt(1350)
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)

        #BACKGROUND
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(21, 21, 30)
        
        #TRANSPARANT LAYER
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(40), top=Pt(130), width=Pt(500), height=Pt(166))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color[0], team_color[1], team_color[2])
        _set_shape_transparency(shape,15000)
        
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(540), top=Pt(130), width=Pt(500), height=Pt(166))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color_2[0], team_color_2[1], team_color_2[2])
        _set_shape_transparency(shape,15000)
        
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(40), top=Pt(298), width=Pt(231), height=Pt(390))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color[0], team_color[1], team_color[2])
        _set_shape_transparency(shape,15000)
        
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(809), top=Pt(298), width=Pt(231), height=Pt(390))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color_2[0], team_color_2[1], team_color_2[2])
        _set_shape_transparency(shape,15000)
        
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(40), top=Pt(690), width=Pt(350), height=Pt(84))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color[0], team_color[1], team_color[2])
        _set_shape_transparency(shape,15000)
        
        shapes = slide.shapes
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(690), top=Pt(690), width=Pt(350), height=Pt(84))
        shape.line.fill.background()
        shapeFill = shape.fill
        shapeFill.solid()
        shapeColour = shapeFill.fore_color
        shapeColour.rgb = RGBColor(team_color_2[0], team_color_2[1], team_color_2[2])
        _set_shape_transparency(shape,15000)
        
        #HEADER
        f1_logo = parent_file / 'data/external/team_logos/F1_75_Logo.png'
        pic = slide.shapes.add_picture(str(f1_logo), Pt(40), Pt(54), height=Pt(32), width= Pt(200))

        title = slide.shapes.title
        title.text = race_name
        title.top = Pt(24)
        title.left = Pt(240)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.name = 'Formula1 Display Bold'

        pic = slide.shapes.add_picture(str(team_logo), Pt(820), Pt(25), height= Pt(100), width=Pt(200))
        
        #STRUCTURE
        line1=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(130), Pt(1000), Pt(2))
        line1.line.fill.background()
        fill = line1.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line2=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(296), Pt(1000), Pt(2))
        line2.line.fill.background()
        fill = line2.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line3=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(688), Pt(1000), Pt(2))
        line3.line.fill.background()
        fill = line3.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line4=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(774), Pt(1000), Pt(2))
        line4.line.fill.background()
        fill = line4.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line5=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(539), Pt(130), Pt(2), Pt(166))
        line5.line.fill.background()
        fill = line5.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line6=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(390), Pt(688), Pt(2), Pt(86))
        line6.line.fill.background()
        fill = line6.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        line7=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(690), Pt(688), Pt(2), Pt(86))
        line7.line.fill.background()
        fill = line7.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(244, 244, 244)

        #TRANSPARENT FIGURES
        
        lap_counter = 0
            
        total_lap_width = 110
        fastest_lap = ast.literal_eval(fastest_lap)
        while lap_counter < len(fastest_lap):
            if fastest_lap[lap_counter] == 0:
                team_color_graph = team_color
            elif fastest_lap[lap_counter] == 1:
                team_color_graph = team_color_2
            else :
                team_color_graph=[120,120,120]
            
            figure_width = 856.8
            lap_width = figure_width/len(fastest_lap)
            shapes = slide.shapes
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(total_lap_width), top=Pt(782), width=Pt(lap_width), height=Pt(385))
            shape.line.fill.background()
            shapeFill = shape.fill
            shapeFill.solid()
            shapeColour = shapeFill.fore_color
            shapeColour.rgb = RGBColor(team_color_graph[0], team_color_graph[1], team_color_graph[2])
            _set_shape_transparency(shape,15000)
            total_lap_width += lap_width

            lap_counter += 1
        
        #FIGURES
        pic = slide.shapes.add_picture(image_file=(str(laptime_scatterplot)), left=Pt(249), top=Pt(283), height=Pt(420), width=Pt(575))

        pic = slide.shapes.add_picture(image_file=(str(driver_1_pace)), left=Pt(19), top=Pt(284), height= Pt(390.5), width=Pt(270))

        pic = slide.shapes.add_picture(image_file=(str(driver_2_pace)), left=Pt(790), top= Pt(284), height= Pt(390), width=Pt(270))
        
        pic = slide.shapes.add_picture(image_file=(str(laptime_comp)), left=Pt(-28), top = Pt(722), height= Pt(500), width=Pt(1108))

        pic = slide.shapes.add_picture(image_file=(str(tyre_strategy)), left=Pt(45), top= Pt(1150),height=Pt(125), width=Pt(943))

        #REFERENCES
        txBox = slide.shapes.add_textbox(left=Pt(70), top= Pt(120), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Driver"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(1010), top= Pt(120), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Driver"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(104), top= Pt(205), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Avg Lap Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)

        txBox = slide.shapes.add_textbox(left=Pt(976), top= Pt(205), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Avg Lap Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(335), top= Pt(120), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Race Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)

        txBox = slide.shapes.add_textbox(left=Pt(535), top= Pt(120), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Race Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(335), top= Pt(205), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Gap"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)

        txBox = slide.shapes.add_textbox(left=Pt(535), top= Pt(205), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Gap"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(165), top= Pt(205), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "IQR"
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(16)
        run.font.bold = True
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(720), top= Pt(205), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "IQR"
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(16)
        run.font.bold = True
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(430), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Race Comparaison"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 270

        txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(500), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Race Comparaison"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 90
        
        txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(920), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lap Time Per Lap"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 270

        txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(990), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lap Time Per Lap"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 90
        
        txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(1210), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Tyre Strategy"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 270

        txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(1260), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Tyre Strategy"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        txBox.rotation = 90
        
        txBox = slide.shapes.add_textbox(left=Pt(130), top= Pt(680), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lap Advantage"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(310), top= Pt(680), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Position"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(540), top= Pt(680), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Safety Car Lap"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(770), top= Pt(680), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Position"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(950), top= Pt(680), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Lap Advantage"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(200), top= Pt(1250), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Total Pit Stop Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(1250), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Pit Stops"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(650), top= Pt(1250), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Pit Stops"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)
        
        txBox = slide.shapes.add_textbox(left=Pt(880), top= Pt(1250), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Total Pit Stop Time"
        run.font.name = 'Formula1 Display Bold'
        run.font.size = Pt(16)
        p.font.color.rgb = RGBColor(120, 120, 120)        
        
        #NEUTRAL VARIABLE
                
        txBox = slide.shapes.add_textbox(left=Pt(540), top= Pt(700), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = safety_car_lap
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        #DRIVER 1 VARIABLE

        txBox = slide.shapes.add_textbox(left=Pt(40), top=Pt(140), width=Pt(500), height=Pt(20))
        txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = name_driver_1
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(335), top= Pt(140), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = total_time_driver_1
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(115), top= Pt(225), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = avg_laptime_driver_1
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(165), top= Pt(225), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(iqr_driver_1)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
                    
        txBox = slide.shapes.add_textbox(left=Pt(335), top= Pt(225), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = driver_1_gap
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
                    
        txBox = slide.shapes.add_textbox(left=Pt(130), top= Pt(700), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = fastest_driver_1
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
                    
        txBox = slide.shapes.add_textbox(left=Pt(310), top= Pt(700), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(driver_1_position)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(200), top= Pt(1270), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = total_duration_pit_driver_1
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(1270), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(number_pit_driver_1)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        #DRIVER 2 VARIABLE
        txBox = slide.shapes.add_textbox(left=Pt(745), top= Pt(140), width=Pt(300), height=Pt(20))
        txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = driver_2_name
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(743), top= Pt(225), width=Pt(300), height=Pt(20))
        txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = avg_laptime_driver_2
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(545), top= Pt(140), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = total_time_driver_2
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(720), top= Pt(225), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(iqr_driver_2)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
                    
        txBox = slide.shapes.add_textbox(left=Pt(545), top= Pt(225), width=Pt(200), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = driver_2_gap
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
                    
        txBox = slide.shapes.add_textbox(left=Pt(770), top= Pt(700), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(driver_2_position)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)    
                
        txBox = slide.shapes.add_textbox(left=Pt(950), top= Pt(700), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = fastest_driver_2
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(650), top= Pt(1270), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(number_pit_driver_2)
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        txBox = slide.shapes.add_textbox(left=Pt(880), top= Pt(1270), width=Pt(0), height=Pt(20))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = total_duration_pit_driver_2
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        counter +=1
        
    except:
        xml_slides = prs.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[counter]) 
        print(f'{team} not in {race_session_name}')
    prs.save(report_folder / f'{race_number}_{race_session_name}.pptx')

for file_path in report_folder.iterdir():
    if file_path.is_file() and file_path.suffix.lower() == '.pptx':
        os.chdir(report_folder)

        subprocess.run(
            [
                "/usr/bin/lowriter",
                "--headless",
                "--convert-to",
                "pdf",
                str(file_path.resolve())
            ],
            check=True
        )
        no_suffix_path = file_path.with_suffix('')
        no_suffix_path.mkdir(parents=True, exist_ok=True)
        file_path_pdf = file_path.with_suffix('.pdf')
        
        os.chdir(no_suffix_path)
        images = convert_from_path(file_path_pdf, dpi=72)

        for i in range(len(images)):
            images[i].save(file_path.stem+ str(i) +'.jpeg', 'JPEG')
        file_path.unlink()
        file_path_pdf.unlink()
    

caption =  ( f"🇦🇺 Round {session.event.RoundNumber} - {session.event.EventName} ({session.event.Country})\n"
            f"📍 {session.event.Location}\n"
            f"🏁 {race_session_name} Session Recap\n"
            f"———\n"
            f"Stay tuned for full weekend coverage of the {session.event.OfficialEventName}!\n"
            f"#F1 #Formula1 #F1{session.event.RoundNumber} #F1{session.event.Country.replace(' ', '')} #F1Weekend #F1Quali #F1Race #F1Team")

creds = read_credentials()
cl = Client()
cl.login(creds['username'], creds['password'])

image_folder = parent_file / f"reports/{session.event.RoundNumber}_{session.event.EventName}_{year}/{session.event.RoundNumber}_{race_session_name}"
image_files = []

for report in os.listdir(image_folder):
    image_files.append(os.path.join(image_folder, report))

cl.delay_range = [1, 3]
cl.album_upload(
    paths=image_files,
    caption=caption
    )