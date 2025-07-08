import matplotlib.pyplot as plt
from matplotlib.collections import LineCollection
import matplotlib.colors as mcolors
from pathlib import Path
from instagrapi import Client
import subprocess
from pdf2image import convert_from_path

import numpy as np 
import pandas as pd

from matplotlib import colormaps
from chromato.spaces import convert

import fastf1
import fastf1.plotting
from fastf1 import utils

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

import os
import csv
import sys
from datetime import datetime

from utils.qual.lap_info import *
from utils.qual.plots import *

parent_file = Path(__file__).resolve().parent.parent

plots_dir = os.path.join(parent_file, 'plots')
if plots_dir not in sys.path:
    sys.path.append(plots_dir)
    print(f"Added {plots_dir} to sys.path")

cache_folder = parent_file / 'cache_folder'
cache_folder.mkdir(exist_ok=True)

fastf1.Cache.enable_cache(cache_folder)

fastf1.plotting.setup_mpl(mpl_timedelta_support=True, misc_mpl_mods=False,
                          color_scheme='fastf1')

year = datetime.today().year
race_number = fastf1.get_events_remaining()['RoundNumber'].values[0]
race_session = 'Q'
race_session_name = 'Qualifying'

session = fastf1.get_session(year, race_number, race_session)
session.load()

q1, q2, q3 = session.laps.split_qualifying_sessions()
is_nat = np.isnat(q1['LapTime'])
q1 = q1[~is_nat]
is_nat = np.isnat(q2['LapTime'])
q2 = q2[~is_nat]
is_nat = np.isnat(q3['LapTime'])
q3 = q3[~is_nat]

event_name = session.event.EventName
figures_folder = parent_file / 'figures' / f"{race_number}_{event_name}_{session.event.year}"
report_folder = parent_file / 'reports' / f"{race_number}_{event_name}_{session.event.year}"

report_folder.mkdir(parents=True, exist_ok=True)
figures_folder.mkdir(parents=True, exist_ok=True)

circuit_info = session.get_circuit_info()
teams = fastf1.plotting.list_team_names(session)

def read_credentials(file = parent_file / "data/raw/login.txt"):
    creds = {}
    with open(file) as f:
        for line in f:
            key, value = line.strip().split("=")
            creds[key] = value
    return creds

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_shape_transparency(shape, alpha):
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))
    
def Hex_RGB(ip):
    return tuple(int(ip[i+1:i+3],16) for i in (0, 2, 4))

delta_per_team = {}
for subsession in [q1, q2, q3]:
    if subsession is q1:
        subsession_name = 'Q1'
    elif subsession is q2:
        subsession_name = 'Q2'
    elif subsession is q3:
        subsession_name = 'Q3'
    for idx, team in enumerate(teams):
        team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
        team_color = fastf1.plotting.get_team_color(team, session=session)
        df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
        team_color_2 = df_color.iat[idx,0]
        try:
            delta_per_team = get_delta_per_team(team, team_drivers, subsession, circuit_info, subsession_name, delta_per_team)
            show_driver_quali_dif_per_lap(team, team_drivers, subsession, circuit_info, subsession_name, team_color, team_color_2, figures_folder)
            show_fastest_lap_per_quali_session(team, team_drivers, subsession, circuit_info, subsession_name, team_color, team_color_2, figures_folder)
            show_corner_advantage_per_quali_session(team, team_drivers, subsession, circuit_info, subsession_name, team_color, team_color_2, delta_per_team, figures_folder)
            create_bar_graph_per_driver(team, team_drivers, subsession, subsession_name, team_color, team_color_2, figures_folder)
        except:
            print(f'No data in {subsession_name} for {team}')

delta_per_team = {}
corner_segments = {}
                    
lap_info_per_driver_q1 = [['EventName', 'Name','LapTime', 'Sector1', 'Sector2', 'Sector3', 'SpeedTrap1', 'SpeedTrap2', 'SpeedFL', 'FullThrottle', 'Brake', 'Cornering', 'Name_1','LapTime_1', 'Sector1_1', 'Sector2_1', 'Sector3_1', 'SpeedTrap1_1', 'SpeedTrap2_1', 'SpeedFL_1', 'FullThrottle_1', 'Brake_1', 'Cornering_1', 'Gap', 'Gap1', 'V_min', 'V_max', 'corner_advantage', 'corner_advantage1']]
lap_info_per_driver_q2 = [['EventName', 'Name','LapTime', 'Sector1', 'Sector2', 'Sector3', 'SpeedTrap1', 'SpeedTrap2', 'SpeedFL', 'FullThrottle', 'Brake', 'Cornering', 'Name_1','LapTime_1', 'Sector1_1', 'Sector2_1', 'Sector3_1', 'SpeedTrap1_1', 'SpeedTrap2_1', 'SpeedFL_1', 'FullThrottle_1', 'Brake_1', 'Cornering_1', 'Gap', 'Gap1', 'V_min', 'V_max', 'corner_advantage', 'corner_advantage1']]
lap_info_per_driver_q3 = [['EventName', 'Name','LapTime', 'Sector1', 'Sector2', 'Sector3', 'SpeedTrap1', 'SpeedTrap2', 'SpeedFL', 'FullThrottle', 'Brake', 'Cornering', 'Name_1','LapTime_1', 'Sector1_1', 'Sector2_1', 'Sector3_1', 'SpeedTrap1_1', 'SpeedTrap2_1', 'SpeedFL_1', 'FullThrottle_1', 'Brake_1', 'Cornering_1', 'Gap', 'Gap1', 'V_min', 'V_max', 'corner_advantage', 'corner_advantage1']]

csv_file_path_q1 = f'{race_number}_{race_session}_drivers_info_Q1.csv'
csv_file_path_q2 = f'{race_number}_{race_session}_drivers_info_Q2.csv'
csv_file_path_q3 = f'{race_number}_{race_session}_drivers_info_Q3.csv'   

csv_file_path_delta_per_team = f'{race_number}_{race_session}_drivers_info_delta_per_team.csv'   
csv_file_path_corner_segments = f'{race_number}_{race_session}_drivers_info_corner_segments.csv'   
 

event_name = session.event.EventName

for idx, team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    try:
        if team_drivers[0] in  q1['Driver'].values and team_drivers[1] in  q1['Driver'].values:
            sub_session = 'Q1'
            delta_per_team.update(show_driver_delta(team, team_drivers, q1, circuit_info, delta_per_team, sub_session))
            corner_segments.update(show_corner_segments(team, team_drivers, q1, circuit_info, corner_segments, sub_session))
            driver_info_figma_q1 = create_csv_lap_info(q1, lap_info_per_driver_q1, team, 'Q1', event_name, race_session, delta_per_team, circuit_info, session)

        if team_drivers[0] in  q2['Driver'].values and team_drivers[1] in  q2['Driver'].values:
            sub_session = 'Q2'
            delta_per_team.update(show_driver_delta(team, team_drivers, q2, circuit_info, delta_per_team, sub_session))
            corner_segments.update(show_corner_segments(team ,team_drivers, q2, circuit_info, corner_segments, sub_session))
            driver_info_figma_q2 = create_csv_lap_info(q2, lap_info_per_driver_q2, team, 'Q2', event_name, race_session, delta_per_team, circuit_info, session)
            
        if team_drivers[0] in  q3['Driver'].values and team_drivers[1] in  q3['Driver'].values:
            sub_session = 'Q3'
            delta_per_team.update(show_driver_delta(team, team_drivers, q3, circuit_info, delta_per_team, sub_session))
            corner_segments.update(show_corner_segments(team, team_drivers, q3, circuit_info, corner_segments, sub_session))
            driver_info_figma_q3 = create_csv_lap_info(q3, lap_info_per_driver_q3, team, 'Q3', event_name, race_session, delta_per_team, circuit_info, session)
    except:
        print(f'{team} not found') 
os.chdir(parent_file / 'data/processed/')
with open(csv_file_path_q1, mode='w', newline='') as file:
    writer = csv.writer(file)
    writer.writerows(driver_info_figma_q1)
        
with open(csv_file_path_q2, mode='w', newline='') as file:
    writer = csv.writer(file)
    writer.writerows(driver_info_figma_q2)
        
with open(csv_file_path_q3, mode='w', newline='') as file:
    writer = csv.writer(file)
    writer.writerows(driver_info_figma_q3)

df_delta_per_team = pd.DataFrame(delta_per_team)
df_delta_per_team.to_csv(csv_file_path_delta_per_team, index=False)
        
df_corner_segments = pd.DataFrame(corner_segments)
df_corner_segments.to_csv(csv_file_path_corner_segments, index=False)

path = parent_file / 'data/processed/'

for subsession in [q1, q2, q3]:
    prs = Presentation()
    if subsession is q1:
        subsession_name = 'Q1'
    elif subsession is q2:
        subsession_name = 'Q2'
    elif subsession is q3:
        subsession_name = 'Q3'
    keyword = f'{race_number}_{race_session}_drivers_info_{subsession_name}'
    for fname in os.listdir(path):
        if keyword in fname:
            driver_data = fname
    os.chdir(path)
    arr = np.loadtxt(driver_data, skiprows=1,
                    delimiter=',', dtype=str)
    keyword = f'{race_number}_{race_session}_drivers_info_corner_segments'
    for fname in os.listdir(path):
        if keyword in fname:
            driver_data = fname
    os.chdir(path)
    corner_segments_arr = np.loadtxt(driver_data, skiprows=1,
                    delimiter=',', dtype=str)
    
    keyword = f'{race_number}_{race_session}_drivers_info_delta_per_teams'
    for fname in os.listdir(path):
        if keyword in fname:
            driver_data = fname
    os.chdir(path)
    delta_per_teams_arr = np.loadtxt(driver_data, skiprows=1,
                    delimiter=',', dtype=str)
    subsession = subsession.drop(subsession[subsession.Deleted == True].index)
    counter = 0
    for idx, team in enumerate(teams):
        team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
        team_color = fastf1.plotting.get_team_color(team, session=session)
        df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
        team_color_2 = df_color.iat[idx,0]
        team_color = Hex_RGB(team_color)
        team_color_2 = Hex_RGB(team_color_2)
        
        if team_drivers[0] in  subsession['Driver'].values and team_drivers[1] in  subsession['Driver'].values:
            try:
                race_name = arr[counter][0]
                team_logo = parent_file / f'data/external/team_logos/{team}.png'
                corner_domination = figures_folder / f'{subsession_name}_{team}_corner_domination.png'
                delta_time = figures_folder / f'{subsession_name}_{team}_deltatime.png'
                speed = figures_folder / f'{subsession_name}_{team}_speed.png'
                
                name_driver_1 = arr[counter][1]
                lap_time_driver_1 = arr[counter][2]
                sector_1_driver_1 = arr[counter][3]
                sector_2_driver_1 = arr[counter][4]
                sector_3_driver_1 = arr[counter][5]
                speed_trap_1_driver_1 = arr[counter][6]
                speed_trap_2_driver_1 = arr[counter][7]
                speed_final_line_driver_1 = arr[counter][8]
                gap_driver_1 = arr[counter][23]
                corner_advantage_driver_1 = arr[counter][27
                                                    ]
                throttling_driver_1 = arr[counter][9]
                braking_driver_1 = arr[counter][10]
                cornering_driver_1 = arr[counter][11]
                bar_graph_driver_1 =  figures_folder / f'{subsession_name}_{team}_driver_1_bar_graph.png'

                name_driver_2 = arr[counter][12]
                lap_time_driver_2 = arr[counter][13]
                sector_1_driver_2 = arr[counter][14]
                sector_2_driver_2 = arr[counter][15]
                sector_3_driver_2 = arr[counter][16]
                speed_trap_1_driver_2 = arr[counter][17]
                speed_trap_2_driver_2 = arr[counter][18]
                speed_final_line_driver_2 = arr[counter][19]
                gap_driver_2 = arr[counter][24]
                corner_advantage_driver_2 = arr[counter][28]
                throttling_driver_2 = arr[counter][20]
                braking_driver_2 = arr[counter][21]
                cornering_driver_2 = arr[counter][22]
                bar_graph_driver_2 = figures_folder / f'{subsession_name}_{team}_driver_2_bar_graph.png'
                
                v_min = arr[counter][-4]
                v_max = arr[counter][-3]
                
                prs.slide_width = Pt(1080)
                prs.slide_height = Pt(1350)
                blank_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(blank_slide_layout)

                #BACKGROUND
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(21, 21, 30)
                
                #TRANSPARANT LAYER
                shapes = slide.shapes
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(40), top=Pt(130), width=Pt(300), height=Pt(435))
                shape.line.fill.background()
                shapeFill = shape.fill
                shapeFill.solid()
                shapeColour = shapeFill.fore_color
                shapeColour.rgb = RGBColor(team_color[0], team_color[1], team_color[2])
                _set_shape_transparency(shape,15000)
                
                shapes = slide.shapes
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(340), top=Pt(385), width=Pt(200), height=Pt(180))
                shape.line.fill.background()
                shapeFill = shape.fill
                shapeFill.solid()
                shapeColour = shapeFill.fore_color
                shapeColour.rgb = RGBColor(team_color[0], team_color[1], team_color[2])
                _set_shape_transparency(shape,15000)
                
                shapes = slide.shapes
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(740), top=Pt(130), width=Pt(300), height=Pt(435))
                shape.line.fill.background()
                shapeFill = shape.fill
                shapeFill.solid()
                shapeColour = shapeFill.fore_color
                shapeColour.rgb = RGBColor(team_color_2[0], team_color_2[1], team_color_2[2])
                _set_shape_transparency(shape,15000)
                
                shapes = slide.shapes
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(540), top=Pt(385), width=Pt(200), height=Pt(180))
                shape.line.fill.background()
                shapeFill = shape.fill
                shapeFill.solid()
                shapeColour = shapeFill.fore_color
                shapeColour.rgb = RGBColor(team_color_2[0], team_color_2[1], team_color_2[2])
                _set_shape_transparency(shape,15000)
                
                #HEADER
                f1_logo = parent_file / 'data/external/team_logos/F1_75_Logo.png'
                pic = slide.shapes.add_picture(str(f1_logo), Pt(40), Pt(54), height=Pt(32), width= Pt(200))

                title = slide.shapes.title
                title.text = race_name
                title.top = Pt(24)
                title.left = Pt(195)

                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                title.text_frame.paragraphs[0].font.size = Pt(38)
                title.text_frame.paragraphs[0].font.name = 'Formula1 Display Bold'

                pic = slide.shapes.add_picture(str(team_logo), Pt(820), Pt(25), height= Pt(100), width=Pt(200))
                
                #STRUCTURE
                line1=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(130), Pt(1000), Pt(2))
                line1.line.fill.background()
                fill = line1.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(244, 244, 244)

                line2=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(565), Pt(1000), Pt(2))
                line2.line.fill.background()
                fill = line2.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(244, 244, 244)

                line3=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(785), Pt(1000), Pt(2))
                line3.line.fill.background()
                fill = line3.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(244, 244, 244)

                line4=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(539), Pt(590), Pt(2), Pt(180))
                line4.line.fill.background()
                fill = line4.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(244, 244, 244)

                line5=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(539), Pt(385), Pt(2), Pt(180))
                line5.line.fill.background()
                fill = line5.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(244, 244, 244)

                #TRANSPARENT FIGURES
                transparent_counter = 0
                df_delta = pd.read_csv(parent_file / f'data/processed/{race_number}_{race_session}_drivers_info_delta_per_team.csv')
                df_delta_columns = list(df_delta.columns)
                delta_index = df_delta_columns.index(f'{subsession_name}_{team}')
                
                df_corner = pd.read_csv(parent_file / f'data/processed/{race_number}_{race_session}_drivers_info_corner_segments.csv')
                df_corner_columns = list(df_corner.columns)
                corner_index = df_corner_columns.index(f'{subsession_name}_{team}')
        
                total_lenght_corner_segments = 44
                while transparent_counter < len(df_delta):
                    if df_delta.iloc[transparent_counter, delta_index] > 0:
                        team_color_graph = team_color
                    else:
                        team_color_graph = team_color_2
                    lenght_corner_segments = df_corner.iloc[transparent_counter, corner_index]
                    
                    shapes = slide.shapes
                    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(total_lenght_corner_segments), top=Pt(810), width=Pt(lenght_corner_segments), height=Pt(305))
                    shape.line.fill.background()
                    shapeFill = shape.fill
                    shapeFill.solid()
                    shapeColour = shapeFill.fore_color
                    shapeColour.rgb = RGBColor(team_color_graph[0], team_color_graph[1], team_color_graph[2])
                    _set_shape_transparency(shape,15000)
                    
                    shapes2 = slide.shapes
                    shape2 = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(total_lenght_corner_segments), top=Pt(1150), width=Pt(lenght_corner_segments), height=Pt(175))
                    shape2.line.fill.background()
                    shape2Fill = shape2.fill
                    shape2Fill.solid()
                    shape2Colour = shape2Fill.fore_color
                    shape2Colour.rgb = RGBColor(team_color_graph[0], team_color_graph[1], team_color_graph[2])
                    _set_shape_transparency(shape2,15000)
                    total_lenght_corner_segments += lenght_corner_segments

                    transparent_counter += 1

                #FIGURES
                pic = slide.shapes.add_picture(image_file=(str(corner_domination)), left=Pt(342), top = Pt(75), height= Pt(395), width=Pt(395))

                pic = slide.shapes.add_picture(image_file=(str(bar_graph_driver_1)), left=Pt(0), top=Pt(565), height= Pt(228), width=Pt(543))

                pic = slide.shapes.add_picture(image_file=(str(bar_graph_driver_2)), left=Pt(523), top= Pt(565), height= Pt(228), width=Pt(543))

                pic = slide.shapes.add_picture(image_file=(str(delta_time)), left=Pt(30), top=Pt(1140), height=Pt(195), width=Pt(1018))

                pic = slide.shapes.add_picture(image_file=(str(speed)), left=Pt(30), top= Pt(798),height=Pt(355), width=Pt(1018))

                #REFERENCES
                txBox = slide.shapes.add_textbox(left=Pt(70), top= Pt(120), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Driver"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                
                txBox = slide.shapes.add_textbox(left=Pt(1010), top= Pt(120), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Driver"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(85), top= Pt(205), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Lap Time"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(995), top= Pt(205), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Lap Time"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(290), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 1"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(290), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 1"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 2"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 2"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 3"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Sector 3"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(290), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Trap 1"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(290), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Trap 1"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Trap 2"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Trap 2"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Final Line"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Speed Final Line"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Gap"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(640), top= Pt(375), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Gap"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Corner Advantage"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(640), top= Pt(460), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Corner Advantage"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(646), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "%LAP TIME"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 270

                txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(685), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "%LAP TIME"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 90

                txBox = slide.shapes.add_textbox(left=Pt(140), top= Pt(550), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "FULL THROTTLE"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(940), top= Pt(550), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "FULL THROTTLE"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(145), top= Pt(615), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "HEAVY BREAKING"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(940), top= Pt(615), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "HEAVY BREAKING"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(120), top= Pt(685), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "CORNERING"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(960), top= Pt(685), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "CORNERING"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)

                txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(935), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "SPEED"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 270

                txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(960), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "SPEED"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 90

                txBox = slide.shapes.add_textbox(left=Pt(0), top= Pt(1195), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "DELTA TIME"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 270

                txBox = slide.shapes.add_textbox(left=Pt(1080), top= Pt(1235), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "DELTA TIME"
                run.font.name = 'Formula1 Display Bold'
                run.font.size = Pt(16)
                p.font.color.rgb = RGBColor(120, 120, 120)
                txBox.rotation = 90

                #DRIVER 1 VARIABLE

                txBox = slide.shapes.add_textbox(left=Pt(40), top=Pt(140), width=Pt(500), height=Pt(20))
                txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = name_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                
                txBox = slide.shapes.add_textbox(left=Pt(125), top= Pt(225), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = lap_time_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(310), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_1_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_2_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)  
                    
                txBox = slide.shapes.add_textbox(left=Pt(103), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_3_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)  
                
                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(310), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_trap_1_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_trap_2_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(260), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_final_line_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = gap_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(430), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = corner_advantage_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(500), top= Pt(572), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = throttling_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) 
                
                txBox = slide.shapes.add_textbox(left=Pt(500), top= Pt(640), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = braking_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                
                txBox = slide.shapes.add_textbox(left=Pt(500), top= Pt(705), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = cornering_driver_1
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                
                #DRIVER 2 VARIABLE
                txBox = slide.shapes.add_textbox(left=Pt(745), top= Pt(140), width=Pt(300), height=Pt(20))
                txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = name_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(950), top= Pt(225), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = lap_time_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(310), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_1_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_2_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(975), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = sector_3_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(310), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_trap_1_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_trap_2_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(817), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = speed_final_line_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(640), top= Pt(395), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = gap_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(640), top= Pt(480), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = corner_advantage_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(580), top= Pt(572), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = throttling_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(580), top= Pt(640), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = braking_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

                txBox = slide.shapes.add_textbox(left=Pt(580), top= Pt(705), width=Pt(0), height=Pt(20))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = cornering_driver_2
                run.font.name = 'Formula1 Display Regular'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                counter +=1
            except:
                xml_slides = prs.slides._sldIdLst  
                slides = list(xml_slides)
                xml_slides.remove(slides[counter]) 
                print(f'{team} not in {subsession_name}')
        
    prs.save(parent_file/ report_folder / f'{race_number}_{subsession_name}_{race_session_name}.pptx')

for file_path in report_folder.iterdir():
    if file_path.is_file() and file_path.suffix.lower() == '.pptx':
        os.chdir(report_folder)

        subprocess.run(
            [
                "/usr/bin/lowriter",
                "--headless",
                "--convert-to",
                "pdf",
                str(file_path.resolve())
            ],
            check=True
        )
        no_suffix_path = file_path.with_suffix('')
        no_suffix_path.mkdir(parents=True, exist_ok=True)
        file_path_pdf = file_path.with_suffix('.pdf')
        
        os.chdir(no_suffix_path)
        images = convert_from_path(file_path_pdf, dpi=72)
       
        for i in range(len(images)):
            images[i].save(file_path.stem + str(i) +'.jpeg', 'JPEG')
        file_path.unlink()
        file_path_pdf.unlink()


caption =  ( f"🇦🇺 Round {session.event.RoundNumber} - {session.event.EventName} ({session.event.Country})\n"
            f"📍 {session.event.Location}\n"
            f"🏁 {race_session_name} Session Recap\n"
            f"———\n"
            f"Stay tuned for full weekend coverage of the {session.event.OfficialEventName}!\n"
            f"#F1 #Formula1 #F1{session.event.RoundNumber} #F1{session.event.Country.replace(' ', '')} #F1Weekend #F1Quali #F1Race #F1Team")

creds = read_credentials()
cl = Client()
cl.login(creds['username'], creds['password'])

for i in range(3):
    image_folder = parent_file / f"reports/{session.event.RoundNumber}_{session.event.EventName}_{year}/{session.event.RoundNumber}_Q{str(i+1)}_{race_session_name}"
    image_files = []
    for report in os.listdir(image_folder):
        image_files.append(os.path.join(image_folder, report))
    cl.delay_range = [1, 3]
    cl.album_upload(
        paths=image_files,
        caption=caption
        )
    cl.delay_range = [1, 3]