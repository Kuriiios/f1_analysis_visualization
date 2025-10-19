from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE


def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element


def _set_shape_transparency(shape, alpha):
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))

def create_line(slide, left:int, top:int, width:int, height:int, red:int, green:int, blue:int):
    line=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, left=Pt(left), top= Pt(top), width=Pt(width), height=Pt(height))
    line.line.fill.background()
    fill = line.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(red, green, blue)

def create_text(slide, left:int, top:int, width:int, height:int, text:str, font_name:str, font_size:int, bold:bool, red:int, green:int, blue:int, align:str, rotation:int):
    txBox = slide.shapes.add_textbox(left=Pt(left), top= Pt(top), width=Pt(width), height=Pt(height))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    match align:
        case "left":
            p.alignment = PP_ALIGN.LEFT
        case "right":
            p.alignment = PP_ALIGN.RIGHT
        case "center":
            p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    p.font.color.rgb = RGBColor(red, green, blue)
    txBox.rotation = rotation

def create_transparant_layer(slide, left:int, top:int, width:int, height:int, red:int, green:int, blue:int):
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left=Pt(left), top=Pt(top), width=Pt(width), height=Pt(height))
    shape.line.fill.background()
    shapeFill = shape.fill
    shapeFill.solid()
    shapeColour = shapeFill.fore_color
    shapeColour.rgb = RGBColor(red, green, blue)
    _set_shape_transparency(shape,15000)