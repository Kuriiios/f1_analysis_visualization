import subprocess
from pdf2image import convert_from_path

from utils.utils import *

import fastf1
from fastf1.ergast import Ergast

import os
import csv

import pandas as pd
import numpy as np
import ast

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ergast = Ergast()

from utils.race.lap_info import *
from utils.race.plots import *

parent_file = Path(__file__).resolve().parent.parent

plots_dir = os.path.join(parent_file, 'plots')
if plots_dir not in sys.path:
    sys.path.append(plots_dir)
    print(f"Added {plots_dir} to sys.path")

cache_folder = parent_file / 'cache_folder'
cache_folder.mkdir(exist_ok=True)

fastf1.Cache.enable_cache(cache_folder)

fastf1.plotting.setup_mpl(mpl_timedelta_support=True, misc_mpl_mods=False,
                          color_scheme='fastf1')

year = int(input('Year ? '))
race_number = int(input('Race Number ? (1-24) '))
race_session = input('Session ? (R, S) ')

session = fastf1.get_session(year, race_number, race_session)
session.load()

if race_session == 'R':
    race_session_name = 'Race'
elif race_session == 'S':
    race_session_name = 'Sprint'

event_name = session.event.EventName   
figures_folder = parent_file / 'figures' / f'{race_number}_{event_name}_{session.event.year}/'
report_folder = parent_file / 'reports' / f"{race_number}_{event_name}_{session.event.year}/"

figures_folder.mkdir(parents=True, exist_ok=True)
report_folder.mkdir(parents=True, exist_ok=True)

pit = ergast.get_pit_stops(season = year, round = race_number )
circuit_info = session.get_circuit_info()
teams = fastf1.plotting.list_team_names(session)

def session_type(race_session):
    match race_session:
        case 'R':
            return 'Race'
        case 'S':
            return 'Sprint'

def read_credentials(file = parent_file / "data/raw/login.txt"):
    creds = {}
    with open(file) as f:
        for line in f:
            key, value = line.strip().split("=")
            creds[key] = value
    return creds

def Hex_RGB(ip):
    return tuple(int(ip[i+1:i+3],16) for i in (0, 2, 4))

def to_str(var):
    return str(list(np.reshape(np.asarray(var), (1, np.size(var)))[0]))[1:-1]

nat_as_integer = np.datetime64('NAT').view('i8')

def isnat(your_datetime):
    dtype_string = str(your_datetime.dtype)
    if 'datetime64' in dtype_string or 'timedelta64' in dtype_string:
        return your_datetime.view('i8') == nat_as_integer
    return False 

def correct_fastest_lap(team, session):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    
    driver_1_laps = session.laps.pick_drivers(team_drivers[0]).pick_laps(range(0, (int(max(session.laps['LapNumber'])) + 1))).reset_index()
    driver_2_laps = session.laps.pick_drivers(team_drivers[1]).pick_laps(range(0, (int(max(session.laps['LapNumber'])) + 1))).reset_index()
    if not driver_1_laps.empty:
        laptime_counter_driver_1 = 0
        for lap in driver_1_laps['LapTime']:
            try:
                if 'NaT' in str(lap):
                    driver_1_laps.loc[laptime_counter_driver_1, 'LapTime'] = driver_1_laps.loc[laptime_counter_driver_1 +1, 'LapStartTime'] - driver_1_laps.loc[laptime_counter_driver_1, 'LapStartTime']
                laptime_counter_driver_1 +=1
            except KeyError:
                break

    laptime_counter_driver_2 = 0
    for lap in driver_2_laps['LapTime']:
        try:
            if 'NaT' in str(lap):
                driver_2_laps.loc[laptime_counter_driver_2, 'LapTime'] = driver_2_laps.loc[laptime_counter_driver_2 +1, 'LapStartTime'] - driver_2_laps.loc[laptime_counter_driver_2, 'LapStartTime']
            laptime_counter_driver_2 +=1
        except KeyError:
            break
    return driver_1_laps, driver_2_laps

for idx,team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    team_color = fastf1.plotting.get_team_color(team, session=session)
    df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
    team_color_2 = df_color.iat[idx,0]
    
    show_pace_comp(team, session, team_color, team_color_2, figures_folder)
    show_laptime_scatterplot(team, team_drivers, session, team_color, team_color_2, figures_folder)
    show_laptime_comp(team, session, team_color, team_color_2, figures_folder)
    show_tyre_strategy(team, team_drivers, session, figures_folder)
    

laptime_scatterplot = {}
laptime_comp = {}
tyre_strategy = {}
pace_comp0 = {}
pace_comp1 = {}

csv_file_path = f'{race_number}_{race_session}_race_info.csv'
csv_file_path_fastest_driver = f'{race_number}_{race_session}_race_fastest_driver.csv'

lap_info_per_team = [['EventName', 'driver_1_name', 'driver_1_position', 'driver_1_gap', 'total_time_driver_1', 'avg_laptime_driver_1','iqr_driver_1', 'number_pit_driver_1', 'total_duration_pit_driver_1', 'driver_2_name', 'driver_2_position', 'driver_2_gap', 'total_time_driver_2', 'avg_laptime_driver_2','iqr_driver_2', 'number_pit_driver_2', 'total_duration_pit_driver_2', 'fastest_driver_1', 'fastest_driver_2', 'safety_car_lap', 'fastest_lap']]
race_info = []
fastest_driver_per_lap_dict = {}
for idx,team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    drivers_info = get_drivers_info(session, team, team_drivers, race_session, year, race_number)
    fastest_driver_per_lap_per_team = get_faster_driver_per_lap(session, team, team_drivers)
    fastest_driver_per_lap_dict.update({f'{team}':fastest_driver_per_lap_per_team})
    team_info = get_lap_repartition( fastest_driver_per_lap_per_team)
    race_info = create_csv_race_info(lap_info_per_team, drivers_info, team_info, event_name, race_session)
os.chdir(parent_file / 'data/processed/')
with open(csv_file_path, mode='w', newline='') as file:
    writer = csv.writer(file)
    writer.writerows(race_info)

keyword = f'{race_number}_{race_session}_race_info'
for fname in os.listdir(parent_file / 'data/processed/'):
    if keyword in fname:
        driver_data = fname
os.chdir(parent_file / 'data/processed/')
arr = pd.read_csv(driver_data)

counter = 0
prs = Presentation()

for idx, team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=session)
    team_color = fastf1.plotting.get_team_color(team, session=session)
    df_color=pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
    team_color_2 = df_color.iat[idx,0]
    team_color = Hex_RGB(team_color)
    team_color_2 = Hex_RGB(team_color_2)
    
    try:
        race_name = arr.iloc[counter]['EventName']
        team_logo = parent_file / f'data/external/team_logos/{team}.png'
        laptime_comp = figures_folder / f'{race_session_name}_{team}_laptime_comp.png'
        laptime_scatterplot = figures_folder / f'{race_session_name}_{team}_laptime_scatterplot.png'
        tyre_strategy = figures_folder / f'{race_session_name}_{team}_tyre_strategy.png'
        driver_1_pace = figures_folder / f'{race_session_name}_{team}_driver_1_pace.png'
        driver_2_pace = figures_folder / f'{race_session_name}_{team}_driver_2_pace.png'
        
        name_driver_1 = arr.iloc[counter]['driver_1_name']
        driver_1_position = str(arr.iloc[counter]['driver_1_position'])
        driver_1_gap = arr.iloc[counter]['driver_1_gap']
        total_time_driver_1 = arr.iloc[counter]['total_time_driver_1']
        avg_laptime_driver_1 = arr.iloc[counter]['avg_laptime_driver_1']
        iqr_driver_1 = str(arr.iloc[counter]['iqr_driver_1'])
        number_pit_driver_1 = str(arr.iloc[counter]['number_pit_driver_1'])
        total_duration_pit_driver_1 = arr.iloc[counter]['total_duration_pit_driver_1']
        
        driver_2_name = arr.iloc[counter]['driver_2_name']
        driver_2_position = str(arr.iloc[counter]['driver_2_position'])
        driver_2_gap = arr.iloc[counter]['driver_2_gap']
        total_time_driver_2 = arr.iloc[counter]['total_time_driver_2']
        avg_laptime_driver_2 = arr.iloc[counter]['avg_laptime_driver_2']
        iqr_driver_2 = str(arr.iloc[counter]['iqr_driver_2'])
        number_pit_driver_2 = str(arr.iloc[counter]['number_pit_driver_2'])
        total_duration_pit_driver_2 = arr.iloc[counter]['total_duration_pit_driver_2']
        
        fastest_driver_1 = arr.iloc[counter]['fastest_driver_1']
        fastest_driver_2 = arr.iloc[counter]['fastest_driver_2']
        safety_car_lap = arr.iloc[counter]['safety_car_lap']
        fastest_lap = arr.iloc[counter]['fastest_lap']
        
        prs.slide_width = Pt(1080)
        prs.slide_height = Pt(1350)
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)

        #BACKGROUND
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(21, 21, 30)
        
        #TRANSPARANT LAYER
        create_transparant_layer(slide, left=40, top=130, width=500, height=166, red=team_color[0], green=team_color[1], blue=team_color[2])
        create_transparant_layer(slide, left=40, top=298, width=231, height=390, red=team_color[0], green=team_color[1], blue=team_color[2])
        create_transparant_layer(slide, left=40, top=690, width=350, height=84, red=team_color[0], green=team_color[1], blue=team_color[2])
        create_transparant_layer(slide, left=540, top=130, width=500, height=166, red=team_color_2[0], green=team_color_2[1], blue=team_color_2[2])
        create_transparant_layer(slide, left=809, top=298, width=231, height=390, red=team_color_2[0], green=team_color_2[1], blue=team_color_2[2])
        create_transparant_layer(slide, left=690, top=690, width=350, height=84, red=team_color_2[0], green=team_color_2[1], blue=team_color_2[2])
        
        #HEADER
        f1_logo = parent_file / 'data/external/team_logos/F1_75_Logo.png'
        pic = slide.shapes.add_picture(str(f1_logo), Pt(40), Pt(54), height=Pt(32), width= Pt(200))

        title = slide.shapes.title
        title.text = race_name
        title.top = Pt(24)
        title.left = Pt(240)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.name = 'Formula1 Display Bold'

        pic = slide.shapes.add_picture(str(team_logo), Pt(820), Pt(25), height= Pt(100), width=Pt(200))
        
        #STRUCTURE
        create_line(slide, left=40, top=130, width=1000, height=2, red=244, green=244, blue=244)
        create_line(slide, left=40, top=296, width=1000, height=2, red=244, green=244, blue=244)
        create_line(slide, left=40, top=688, width=1000, height=2, red=244, green=244, blue=244)
        create_line(slide, left=40, top=774, width=1000, height=2, red=244, green=244, blue=244)
        create_line(slide, left=539, top=130, width=2, height=166, red=244, green=244, blue=244)
        create_line(slide, left=390, top=688, width=2, height=86, red=244, green=244, blue=244)
        create_line(slide, left=690, top=688, width=2, height=86, red=244, green=244, blue=244)

        #TRANSPARENT FIGURES
        lap_counter = 0
            
        total_lap_width = 110
        fastest_lap = ast.literal_eval(fastest_lap)
        while lap_counter < len(fastest_lap):
            if fastest_lap[lap_counter] == 0:
                team_color_graph = team_color
            elif fastest_lap[lap_counter] == 1:
                team_color_graph = team_color_2
            else :
                team_color_graph=[120,120,120]
            
            figure_width = 856.8
            lap_width = figure_width/len(fastest_lap)
            create_transparant_layer(slide, left=total_lap_width, top=782, width=lap_width, height=385, red=team_color_graph[0], green=team_color_graph[1], blue=team_color_graph[2])
            total_lap_width += lap_width
            lap_counter += 1
        
        #FIGURES
        pic = slide.shapes.add_picture(image_file=(str(laptime_scatterplot)), left=Pt(249), top=Pt(283), height=Pt(420), width=Pt(575))
        pic = slide.shapes.add_picture(image_file=(str(driver_1_pace)), left=Pt(19), top=Pt(284), height= Pt(390.5), width=Pt(270))
        pic = slide.shapes.add_picture(image_file=(str(driver_2_pace)), left=Pt(790), top= Pt(284), height= Pt(390), width=Pt(270))
        pic = slide.shapes.add_picture(image_file=(str(laptime_comp)), left=Pt(-28), top = Pt(722), height= Pt(500), width=Pt(1108))
        pic = slide.shapes.add_picture(image_file=(str(tyre_strategy)), left=Pt(45), top= Pt(1150),height=Pt(125), width=Pt(943))

        #REFERENCES
        create_text(slide, left=540, top=680, width=0, height=20, text="Safety Car Lap", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)

        create_text(slide, left=70, top=120, width=0, height=20, text="Driver", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=104, top=205, width=0, height=20, text="Avg Lap Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=335, top=120, width=200, height=20, text="Race Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=335, top=205, width=200, height=20, text="Gap", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=165, top=205, width=200, height=20, text="IQR", font_name="Formula1 Display Bold", font_size=16, bold=True, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=0, top=430, width=0, height=20, text="Race Comparaison", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=270)
        create_text(slide, left=0, top=920, width=0, height=20, text="Lap Time Per Lap", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=270)
        create_text(slide, left=0, top=1210, width=0, height=20, text="Tyre Strategy", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=270)
        create_text(slide, left=130, top=680, width=0, height=20, text="Lap Advantage", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=310, top=680, width=0, height=20, text="Position", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=430, top=1250, width=0, height=20, text="Pit Stops", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=200, top=1250, width=0, height=20, text="Total Pit Stop Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)

        create_text(slide, left=1010, top=120, width=0, height=20, text="Driver", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=976, top=205, width=0, height=20, text="Avg Lap Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=535, top=120, width=200, height=20, text="Race Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=720, top=205, width=200, height=20, text="IQR", font_name="Formula1 Display Bold", font_size=16, bold=True, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=1080, top=500, width=0, height=20, text="Race Comparaison", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=90)
        create_text(slide, left=1080, top=990, width=0, height=20, text="Lap Time Per Lap", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=90)
        create_text(slide, left=1080, top=1260, width=0, height=20, text="Tyre Strategy", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=90)
        create_text(slide, left=950, top=680, width=0, height=20, text="Lap Advantage", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=770, top=680, width=0, height=20, text="Position", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=650, top=1250, width=0, height=20, text="Pit Stops", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)
        create_text(slide, left=880, top=1250, width=0, height=20, text="Total Pit Stop Time", font_name="Formula1 Display Bold", font_size=16, bold=False, red=120, green=120, blue=120, align="right", rotation=0)      
        
        #NEUTRAL VARIABLE
        create_text(slide, left=540, top=700, width=0, height=20, text=safety_car_lap, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        
        #DRIVER 1 VARIABLE
        create_text(slide, left=40, top=140, width=300, height=20, text=name_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=335, top=140, width=200, height=20, text=total_time_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=115, top=225, width=0, height=20, text=avg_laptime_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=165, top=225, width=200, height=20, text=iqr_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=335, top=225, width=200, height=20, text=driver_1_gap, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=130, top=700, width=0, height=20, text=fastest_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=310, top=700, width=0, height=20, text=driver_1_position, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=200, top=1270, width=0, height=20, text=total_duration_pit_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        create_text(slide, left=430, top=1270, width=0, height=20, text=number_pit_driver_1, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
        
        #DRIVER 2 VARIABLE
        create_text(slide, left=745, top=140, width=300, height=20, text=driver_2_name, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=815, top=225, width=300, height=20, text=avg_laptime_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=545, top=140, width=200, height=20, text=total_time_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=720, top=225, width=200, height=20, text=iqr_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=545, top=225, width=200, height=20, text=driver_2_gap, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=770, top=700, width=0, height=20, text=driver_2_position, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=950, top=700, width=0, height=20, text=fastest_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=650, top=1270, width=0, height=20, text=number_pit_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
        create_text(slide, left=880, top=1270, width=0, height=20, text=total_duration_pit_driver_2, font_name="Formula1 Display Regular", font_size=30, bold=True, red=255, green=255, blue=255, align="right", rotation=0)
      
        counter +=1
        
    except:
        xml_slides = prs.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[counter]) 
        print(f'{team} not in {race_session_name}')
    prs.save(report_folder / f'{race_number}_{race_session_name}.pptx')

for file_path in report_folder.iterdir():
    if file_path.is_file() and file_path.suffix.lower() == '.pptx':
        os.chdir(report_folder)

        subprocess.run(
            [
                "/usr/bin/lowriter",
                "--headless",
                "--convert-to",
                "pdf",
                str(file_path.resolve())
            ],
            check=True
        )
        no_suffix_path = file_path.with_suffix('')
        no_suffix_path.mkdir(parents=True, exist_ok=True)
        file_path_pdf = file_path.with_suffix('.pdf')
        
        os.chdir(no_suffix_path)
        images = convert_from_path(file_path_pdf, dpi=72)

        for i in range(len(images)):
            images[i].save(file_path.stem+ str(i) +'.jpeg', 'JPEG')
        file_path.unlink()
        file_path_pdf.unlink()
