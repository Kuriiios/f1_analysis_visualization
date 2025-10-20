import pandas as pd
import os  
from datetime import timedelta, datetime
from pathlib import Path
import subprocess

from utils.utils import *

from fastf1.ergast import Ergast
import fastf1.plotting

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

from pdf2image import convert_from_path

ergast = Ergast()

parent_file = Path(__file__).resolve().parent.parent

year = int(input("Year ? "))
race_from = int(input("From ? "))
race_to = int(input("To ? "))

report_folder = parent_file / 'reports'

def parse_laptime_str(lap_string):
    """Parses a lap time string (MM:SS.fff) into a timedelta object."""
    if isinstance(lap_string, str) and ':' in lap_string and '.' in lap_string:
        minutes, rest = lap_string.split(':')
        seconds, milliseconds = rest.split('.')
        return timedelta(minutes=int(minutes), seconds=int(seconds), milliseconds=int(milliseconds))
    else:
        return timedelta(0)

def Hex_RGB(ip):
    return tuple(int(ip[i+1:i+3],16) for i in (0, 2, 4))

race_years = list(range(race_from, race_to+1))
recaps = {}
zero_time = timedelta(0)
threshold_multiplier = 1.5

qual_sessions = range(1, 4)


for race_number in race_years:
    for race_session in ('S', 'R', 'SQ', 'Q'):
        try:
            session = fastf1.get_session(year, race_number, race_session)
            session.load()
            print(session)
            teams = fastf1.plotting.list_team_names(session)
        except:
            continue
        if race_session == 'S' or race_session == 'R':
            keyword = f'{race_number}_{race_session}_race_info'
            driver_data_file = None
            for fname in os.listdir(parent_file / 'data/processed/'):
                if keyword in fname:
                    driver_data_file = fname
                    break

            if driver_data_file:
                os.chdir(parent_file / 'data/processed/')
                arr = pd.read_csv(driver_data_file)
                num_rows = len(arr)

                temp_recap = {}
                print(arr)
                for i in range(num_rows):
                    
                    name_driver_1 = arr['driver_1_name'][i]
                    name_driver_2 = arr['driver_2_name'][i]

                    name = f'{name_driver_1}_{name_driver_2}'
                    iqr_driver_1 = float(arr['iqr_driver_1'][i][:-2])
                    iqr_driver_2 = float(arr['iqr_driver_2'][i][:-2])

                    laptime_driver_1_str = arr['avg_laptime_driver_1'][i]
                    laptime_driver_2_str = arr['avg_laptime_driver_2'][i]

                    fastest_lap_driver_1_str = arr['fastest_driver_1'][i]
                    fastest_lap_driver_2_str = arr['fastest_driver_2'][i]
                    safety_car_lap_str = arr['safety_car_lap'][i]

                    avg_seconds_driver_1 = parse_laptime_str(laptime_driver_1_str)
                    avg_seconds_driver_2 = parse_laptime_str(laptime_driver_2_str)

                    if (iqr_driver_1 != 0 and iqr_driver_2 != 0 and avg_seconds_driver_1 != zero_time and avg_seconds_driver_2 != zero_time
                            and avg_seconds_driver_1 <= threshold_multiplier * avg_seconds_driver_2 and avg_seconds_driver_2 <= threshold_multiplier * avg_seconds_driver_1):

                        if datetime.strptime(laptime_driver_1_str, '%M:%S.%f') < datetime.strptime(laptime_driver_2_str, '%M:%S.%f'):
                            gap_driver_1 = -round((pd.to_timedelta(datetime.strptime(laptime_driver_2_str, '%M:%S.%f') - datetime.strptime(laptime_driver_1_str, '%M:%S.%f')).total_seconds()), 3)
                        else:
                            gap_driver_1 = round((pd.to_timedelta(datetime.strptime(laptime_driver_1_str, '%M:%S.%f') - datetime.strptime(laptime_driver_2_str, '%M:%S.%f')).total_seconds()), 3)

                        gap_driver_2 = -gap_driver_1

                        lap_advantage_1 = round((int(fastest_lap_driver_1_str[:-3])/(int(fastest_lap_driver_1_str[-2:])-int(safety_car_lap_str[:-3])))*100)
                        lap_advantage_2 = round((int(fastest_lap_driver_2_str[:-3])/(int(fastest_lap_driver_2_str[-2:])-int(safety_car_lap_str[:-3])))*100)

                        pit_stop_number_driver_1 = int(arr['number_pit_driver_1'][i])
                        pit_stop_number_driver_2 = int(arr['number_pit_driver_2'][i])

                        pit_stop_time_driver_1_td = parse_laptime_str(arr['total_duration_pit_driver_1'][i][:-2])
                        pit_stop_time_driver_2_td = parse_laptime_str(arr['total_duration_pit_driver_2'][i][:-2])

                        temp_recap.setdefault(name, {'first_race': 0, 'last_race': 0, 'total_iqr_driver_1': 0, 'total_gap_driver_1': 0, 'total_lap_advantage_driver_1': 0, 'total_pit_stop_time_driver_1': timedelta(0), 'total_pit_stop_number_driver_1': 0, 'count_driver_1': 0, 'total_pit_stop_number Driver 1': 0, 'total_qual_gap Driver 1': 0, 'total_corner_advantage Driver 1': 0, 'qual_count Driver 1': 0, 
                                                    'total_iqr_driver_2': 0, 'total_gap_driver_2': 0, 'total_lap_advantage_driver_2': 0, 'total_pit_stop_time_driver_2': timedelta(0), 'total_pit_stop_number_driver_2': 0, 'count_driver_2': 0, 'total_pit_stop_number Driver 2': 0, 'total_qual_gap Driver 2': 0, 'total_corner_advantage Driver 2': 0, 'qual_count Driver 2': 0})
                        
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['first_race'] = min(race_number, temp_recap[f'{name_driver_1}_{name_driver_2}']['first_race'])
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['last_race'] = max(race_number, temp_recap[f'{name_driver_1}_{name_driver_2}']['last_race'])

                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_iqr_driver_1'] += iqr_driver_1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_gap_driver_1'] += gap_driver_1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_lap_advantage_driver_1'] += lap_advantage_1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_pit_stop_time_driver_1'] += pit_stop_time_driver_1_td
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_pit_stop_number_driver_1'] += pit_stop_number_driver_1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['count_driver_1'] += 1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_qual_gap Driver 1'] = 0
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_corner_advantage Driver 1'] = 0

                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_iqr_driver_2'] += iqr_driver_2
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_gap_driver_2'] += gap_driver_2
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_lap_advantage_driver_2'] += lap_advantage_2
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_pit_stop_time_driver_2'] += pit_stop_time_driver_2_td
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_pit_stop_number_driver_2'] += pit_stop_number_driver_2
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['count_driver_2'] += 1
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_qual_gap Driver 2'] = 0
                        temp_recap[f'{name_driver_1}_{name_driver_2}']['total_corner_advantage Driver 2'] = 0 

                for name, data in temp_recap.items():

                    recaps.setdefault(name, {'First Race': 0, 'Last Race': 0, 'Average IQR Driver 1': 0, 'Average Gap Race Lap Driver 1': 0, 'Lap Advantage Driver 1': 0, 'Average Pit Stop Time Driver 1': timedelta(0), 'total_count Driver 1': 0, 'total_pit_stop_number Driver 1': 0, 'total_qual_gap Driver 1': 0, 'total_corner_advantage Driver 1': 0, 'qual_count Driver 1': 0, 
                                            'Average IQR Driver 2': 0, 'Average Gap Race Lap Driver 2': 0, 'Lap Advantage Driver 2': 0, 'Average Pit Stop Time Driver 2': timedelta(0), 'total_count Driver 2': 0, 'total_pit_stop_number Driver 2': 0, 'total_qual_gap Driver 2': 0, 'total_corner_advantage Driver 2': 0, 'qual_count Driver 2': 0})
                    
                    recaps[name]['First Race'] = data['first_race'] + 1
                    recaps[name]['Last Race'] = data['last_race']

                    recaps[name]['Average IQR Driver 1'] += data['total_iqr_driver_1']
                    recaps[name]['Average Gap Race Lap Driver 1'] += data['total_gap_driver_1']
                    recaps[name]['Lap Advantage Driver 1'] += data['total_lap_advantage_driver_1']
                    recaps[name]['Average Pit Stop Time Driver 1'] += data['total_pit_stop_time_driver_1']
                    recaps[name]['total_count Driver 1'] += data['count_driver_1']
                    recaps[name]['total_pit_stop_number Driver 1'] += data['total_pit_stop_number_driver_1']
                    
                    recaps[name]['Average IQR Driver 2'] += data['total_iqr_driver_2']
                    recaps[name]['Average Gap Race Lap Driver 2'] += data['total_gap_driver_2']
                    recaps[name]['Lap Advantage Driver 2'] += data['total_lap_advantage_driver_2']
                    recaps[name]['Average Pit Stop Time Driver 2'] += data['total_pit_stop_time_driver_2']
                    recaps[name]['total_count Driver 2'] += data['count_driver_2']
                    recaps[name]['total_pit_stop_number Driver 2'] += data['total_pit_stop_number_driver_2']
        
        if race_session == 'SQ' or race_session == 'Q':            
            for qual_session in qual_sessions:
                keyword = f'{race_number}_{race_session}_drivers_info_Q{qual_session}'
                driver_data_file = None
                for fname in os.listdir(parent_file / 'data/processed/'):
                    if keyword in fname:
                        driver_data_file = fname
                        break

                if driver_data_file:
                    os.chdir(parent_file / 'data/processed/')
                    arr = pd.read_csv(driver_data_file)

                    num_rows = len(arr)

                    for i in range(num_rows):
                        name_driver_1 = arr['Name'][i]
                        name_driver_2 = arr['Name_1'][i]

                        name = f'{name_driver_1}_{name_driver_2}'

                        laptime_driver_1_str = arr['LapTime'][i]
                        laptime_driver_2_str = arr['LapTime_1'][i]

                        avg_seconds_driver_1 = parse_laptime_str(laptime_driver_1_str)
                        avg_seconds_driver_2 = parse_laptime_str(laptime_driver_2_str)

                        if (avg_seconds_driver_1 != zero_time and avg_seconds_driver_2 != zero_time
                                and avg_seconds_driver_1 <= threshold_multiplier * avg_seconds_driver_2 and avg_seconds_driver_2 <= threshold_multiplier * avg_seconds_driver_1):

                            time_driver_1 = datetime.strptime(laptime_driver_1_str, '%M:%S.%f')
                            time_driver_2 = datetime.strptime(laptime_driver_2_str, '%M:%S.%f')

                            gap = round((pd.to_timedelta(time_driver_1 - time_driver_2).total_seconds()), 3)

                            corner_adv_1_num, corner_adv_1_den = map(int, arr['corner_advantage'][i].split('/'))
                            corner_advantage_driver_1 = (corner_adv_1_num / corner_adv_1_den) * 100

                            corner_adv_2_num, corner_adv_2_den = map(int, arr['corner_advantage1'][i].split('/'))
                            corner_advantage_driver_2 = (corner_adv_2_num / corner_adv_2_den) * 100
                            
                            if not name in temp_recap:
                                temp_recap.setdefault(name, {'Average IQR Driver 1': 0, 'Average Gap Race Lap Driver 1': 0, 'Lap Advantage Driver 1': 0, 'Average Pit Stop Time Driver 1': timedelta(0), 'total_count Driver 1': 0, 'total_pit_stop_number Driver 1': 0, 'total_qual_gap Driver 1': 0, 'total_corner_advantage Driver 1': 0, 'qual_count Driver 1': 0, 
                                                            'Average IQR Driver 2': 0, 'Average Gap Race Lap Driver 2': 0, 'Lap Advantage Driver 2': 0, 'Average Pit Stop Time Driver 2': timedelta(0), 'total_count Driver 2': 0, 'total_pit_stop_number Driver 2': 0, 'total_qual_gap Driver 2': 0, 'total_corner_advantage Driver 2': 0, 'qual_count Driver 2': 0})
                            
                            temp_recap[f'{name_driver_1}_{name_driver_2}']['total_qual_gap Driver 1'] += gap
                            temp_recap[f'{name_driver_1}_{name_driver_2}']['total_corner_advantage Driver 1'] += corner_advantage_driver_1
                            temp_recap[f'{name_driver_1}_{name_driver_2}']['qual_count Driver 1'] += 1

                            temp_recap[f'{name_driver_1}_{name_driver_2}']['total_qual_gap Driver 2'] -= gap
                            temp_recap[f'{name_driver_1}_{name_driver_2}']['total_corner_advantage Driver 2'] += corner_advantage_driver_2
                            temp_recap[f'{name_driver_1}_{name_driver_2}']['qual_count Driver 2'] += 1

                    for name, data in temp_recap.items():
                        
                        if not name in recaps:
                            recaps.setdefault(name, {'Average IQR Driver 1': 0, 'Average Gap Race Lap Driver 1': 0, 'Lap Advantage Driver 1': 0, 'Average Pit Stop Time Driver 1': timedelta(0), 'total_count Driver 1': 0, 'total_pit_stop_number Driver 1': 0, 'total_qual_gap Driver 1': 0, 'total_corner_advantage Driver 1': 0, 'qual_count Driver 1': 0, 
                                                     'Average IQR Driver 2': 0, 'Average Gap Race Lap Driver 2': 0, 'Lap Advantage Driver 2': 0, 'Average Pit Stop Time Driver 2': timedelta(0), 'total_count Driver 2': 0, 'total_pit_stop_number Driver 2': 0, 'total_qual_gap Driver 2': 0, 'total_corner_advantage Driver 2': 0, 'qual_count Driver 2': 0})
                        
                        recaps[name]['total_qual_gap Driver 1'] += data['total_qual_gap Driver 1']
                        recaps[name]['total_corner_advantage Driver 1'] += data['total_corner_advantage Driver 1']
                        recaps[name]['qual_count Driver 1'] += data['qual_count Driver 1']

                        recaps[name]['total_qual_gap Driver 2'] += data['total_qual_gap Driver 2']
                        recaps[name]['total_corner_advantage Driver 2'] += data['total_corner_advantage Driver 2']
                        recaps[name]['qual_count Driver 2'] += data['qual_count Driver 2']


final_race_recaps = {}
for name, data in recaps.items():
    
    first_race = data['First Race']
    last_race = data['Last Race']

    count_driver_1 = data['total_count Driver 1']
    avg_iqr_driver_1 = round(data['Average IQR Driver 1'] / count_driver_1, 3) if count_driver_1 > 0 else "N/A"
    avg_gap_driver_1 = round(data['Average Gap Race Lap Driver 1'] / count_driver_1, 3) if count_driver_1 > 0 else "N/A"
    lap_advantage_driver_1 = round(data['Lap Advantage Driver 1'] / count_driver_1) if count_driver_1 > 0 else "N/A"
    avg_pit_stop_time_driver_1 = round(data['Average Pit Stop Time Driver 1'].total_seconds() / data['total_pit_stop_number Driver 1'], 3) if data['total_pit_stop_number Driver 1'] > 0 else "N/A"

    count_qual_driver_1 = data['qual_count Driver 1']
    avg_qual_gap_driver_1 = round(data['total_qual_gap Driver 1'] / count_qual_driver_1, 3) if count_qual_driver_1 > 0 else "N/A"
    avg_corner_advantage_driver_1 = round(data['total_corner_advantage Driver 1'] / count_qual_driver_1) if count_qual_driver_1 > 0 else "N/A"
    
    count_driver_2 = data['total_count Driver 2']
    avg_iqr_driver_2 = round(data['Average IQR Driver 2'] / count_driver_2, 3) if count_driver_2 > 0 else "N/A"
    avg_gap_driver_2 = round(data['Average Gap Race Lap Driver 2'] / count_driver_2, 3) if count_driver_2 > 0 else "N/A"
    lap_advantage_driver_2 = round(data['Lap Advantage Driver 2'] / count_driver_2) if count_driver_2 > 0 else "N/A"
    avg_pit_stop_time_driver_2 = round(data['Average Pit Stop Time Driver 2'].total_seconds() / data['total_pit_stop_number Driver 2'], 3) if data['total_pit_stop_number Driver 2'] > 0 else "N/A"

    count_qual_driver_2 = data['qual_count Driver 2']
    avg_qual_gap_driver_2 = round(data['total_qual_gap Driver 2'] / count_qual_driver_2, 3) if count_qual_driver_2 > 0 else "N/A"
    avg_corner_advantage_driver_2 = round(data['total_corner_advantage Driver 2'] / count_qual_driver_2) if count_qual_driver_2 > 0 else "N/A"

    final_race_recaps[name] = {
        'First Race' : first_race,
        'Last Race' : last_race,
        'Average Gap Quali Lap Driver 1': f'{avg_qual_gap_driver_1}s',
        'Corner Advantage Driver 1': f'{avg_corner_advantage_driver_1}%',
        'Average IQR Driver 1': f'{avg_iqr_driver_1}s',
        'Average Gap Race Lap Driver 1': f'{avg_gap_driver_1}s',
        'Lap Advantage Driver 1': f'{lap_advantage_driver_1}%',
        'Average Pit Stop Time Driver 1': f'{avg_pit_stop_time_driver_1}s',

        'Average Gap Quali Lap Driver 2': f'{avg_qual_gap_driver_2}s',
        'Corner Advantage Driver 2': f'{avg_corner_advantage_driver_2}%',
        'Average IQR Driver 2': f'{avg_iqr_driver_2}s',
        'Average Gap Race Lap Driver 2': f'{avg_gap_driver_2}s',
        'Lap Advantage Driver 2': f'{lap_advantage_driver_2}%',
        'Average Pit Stop Time Driver 2': f'{avg_pit_stop_time_driver_2}s'
    }

teams_list = list(final_race_recaps)
prs = Presentation()

compared_teams = set()

for race_number in race_years:

    for race_session in ('R', 'S', 'Q', 'SQ'):  # Iterate through sessions
        try :
            session = fastf1.get_session(year, race_number, race_session) # session is already defined in the main loop, but it's needed here as well.
            session.load()
            teams = fastf1.plotting.list_team_names(session)

            for team in teams:
                drivers_name = fastf1.plotting.get_driver_names_by_team(team, session)

                if len(drivers_name) == 2:
                    name_driver_1 = drivers_name[0]
                    name_driver_2 = drivers_name[1]

                    team_name = fastf1.plotting.get_team_name_by_driver(name_driver_1, session)
                    team_tuple = tuple(sorted((name_driver_1, name_driver_2)))

                    if team_tuple not in compared_teams:
                        compared_teams.add(team_tuple)
                        team_logo = parent_file / f'data/external/team_logos/{team_name}.png'
                        
                        df_color = pd.read_csv(parent_file / "data/raw/second_color.csv", index_col='team')
                        team_color = fastf1.plotting.get_team_color(team_name, session=session)
                        team_color_2 = None
                        for idx, row in df_color.iterrows():
                            if row.name == team_name:
                                team_color_2 = row['color']
                                break  

                        team_color = Hex_RGB(team_color)
                        team_color_2 = Hex_RGB(team_color_2)

                        first_race = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('First Race', 'N/A')
                        last_race = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Last Race', 'N/A')

                        average_gap_quali_lap_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Gap Quali Lap Driver 1', 'N/A')
                        corner_advantage_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Corner Advantage Driver 1', 'N/A')
                        average_iqr_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average IQR Driver 1', 'N/A')
                        average_gap_race_lap_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Gap Race Lap Driver 1', 'N/A')
                        lap_advantage_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Lap Advantage Driver 1', 'N/A')
                        average_pit_stop_time_driver_1 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Pit Stop Time Driver 1', 'N/A')

                        average_gap_quali_lap_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Gap Quali Lap Driver 2', 'N/A')
                        corner_advantage_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Corner Advantage Driver 2', 'N/A')
                        average_iqr_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average IQR Driver 2', 'N/A')
                        average_gap_race_lap_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Gap Race Lap Driver 2', 'N/A')
                        lap_advantage_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Lap Advantage Driver 2', 'N/A')
                        average_pit_stop_time_driver_2 = final_race_recaps.get(f'{name_driver_1}_{name_driver_2}', {}).get('Average Pit Stop Time Driver 2', 'N/A')


                        prs.slide_width = Pt(1080)
                        prs.slide_height = Pt(1080)
                        blank_slide_layout = prs.slide_layouts[5]
                        slide = prs.slides.add_slide(blank_slide_layout)

                        #BACKGROUND
                        background = slide.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(21, 21, 30)

                        #HEADER
                        f1_logo = parent_file / 'data/external/team_logos/F1_75_Logo.png'
                        pic = slide.shapes.add_picture(str(f1_logo), Pt(40), Pt(54), height=Pt(32), width= Pt(200))

                        title = slide.shapes.title
                        title.top = Pt(26)
                        title.left = Pt(200)
                        title.text = f'Recap for  from race {first_race} to {last_race}'
                        
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        title.text_frame.paragraphs[0].font.size = Pt(38)
                        title.text_frame.paragraphs[0].font.name = 'Formula1 Display Bold'

                        pic = slide.shapes.add_picture(str(team_logo), Pt(820), Pt(25), height= Pt(100), width=Pt(200))

                        #STRUCTURE
                        create_line(slide, left=40, top=130, width=1000, height=2, red=244, green=244, blue=244)
                        #TRANSPARANT LAYER
                        create_transparant_layer(slide, left=40, top=130, width=500, height=910, red=team_color[0], green=team_color[1], blue=team_color[2])
                        create_transparant_layer(slide, left=540, top=130, width=500, height=910, red=team_color_2[0], green=team_color_2[1], blue=team_color_2[2])

                        #REFERENCES
                        create_text(slide, left=330, top=200, width=400, height=32, text="Average Gap Fastest Quali Lap per Session", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
                        create_text(slide, left=330, top=340, width=400, height=32, text="Corner Advantage Quali Lap", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
                        create_text(slide, left=330, top=480, width=400, height=32, text="Average IQR Race", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
                        create_text(slide, left=330, top=620, width=400, height=32, text="Average Gap for Fast Race Laps", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
                        create_text(slide, left=330, top=760, width=400, height=32, text="Lap Advantage per Race", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
                        create_text(slide, left=330, top=900, width=400, height=32, text="Average Pit Stop Time", font_name="Formula1 Display Bold", font_size=32, bold=False, red=255, green=255, blue=255, align="right", rotation=0)
            
                        #DRIVER 1 VARIABLE
                        create_text(slide, left=60, top=135, width=500, height=20, text=name_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=265, width=500, height=20, text=average_gap_quali_lap_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=405, width=500, height=20, text=corner_advantage_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=545, width=500, height=20, text=average_iqr_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=685, width=500, height=20, text=average_gap_race_lap_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=825, width=500, height=20, text=lap_advantage_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=100, top=965, width=500, height=20, text=average_pit_stop_time_driver_1, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)

                        #DRIVER 2 VARIABLE
                        create_text(slide, left=520, top=135, width=500, height=20, text=name_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=265, width=500, height=20, text=average_gap_quali_lap_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=405, width=500, height=20, text=corner_advantage_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=545, width=500, height=20, text=average_iqr_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=685, width=500, height=20, text=average_gap_race_lap_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=825, width=500, height=20, text=lap_advantage_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)
                        create_text(slide, left=480, top=965, width=500, height=20, text=average_pit_stop_time_driver_2, font_name="Formula1 Display Regular", font_size=40, bold=True, red=255, green=255, blue=255, align="left", rotation=0)

                prs.save(parent_file/ report_folder / f'Recap_race_{race_from}_{race_to}.pptx')
        except:
            continue

for file_path in report_folder.iterdir():
    if file_path.is_file() and file_path.suffix.lower() == '.pptx':
        os.chdir(report_folder)

        subprocess.run(
            [
                "/usr/bin/lowriter",
                "--headless",
                "--convert-to",
                "pdf",
                str(file_path.resolve())
            ],
            check=True
        )
        no_suffix_path = file_path.with_suffix('')
        no_suffix_path.mkdir(parents=True, exist_ok=True)
        file_path_pdf = file_path.with_suffix('.pdf')
        
        os.chdir(no_suffix_path)
        images = convert_from_path(file_path_pdf, dpi=72)
       
        for i in range(len(images)):
            images[i].save(file_path.stem + str(i) +'.jpeg', 'JPEG')
        file_path.unlink()
        file_path_pdf.unlink()
