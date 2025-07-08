from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

from instagrapi import Client
from pathlib import Path
from pdf2image import convert_from_path
import subprocess
import sys

import fastf1
import fastf1.plotting

import pandas as pd
import numpy as np
import os  
from datetime import timedelta
import csv
from fastf1.ergast import Ergast

parent_file = Path(__file__).resolve().parent.parent

plots_dir = os.path.join(parent_file, 'plots')
if plots_dir not in sys.path:
    sys.path.append(plots_dir)
    print(f"Added {plots_dir} to sys.path")

cache_folder = parent_file / 'cache_folder'
cache_folder.mkdir(exist_ok=True)

fastf1.Cache.enable_cache(cache_folder)

fastf1.plotting.setup_mpl(mpl_timedelta_support=True, misc_mpl_mods=False,
                          color_scheme='fastf1')

ergast = Ergast()

year = int(input('Year ? '))
race_number = int(input('Race Number ? (1-24) '))
free_practice_session = 'FP' + (input('Practice Session ? (1-3) '))
post_option = input('Do you want to post it immediatly ? (Y/N) ')

event = fastf1.get_event(year, race_number)

practice_session = fastf1.get_session(year, race_number, free_practice_session)
practice_session.load(weather=True)

teams = fastf1.plotting.list_team_names(practice_session)
event_name = practice_session.event.EventName

report_folder = parent_file / 'reports' / f"{race_number}_{event_name}_{practice_session.event.year}"
report_folder.mkdir(parents=True, exist_ok=True)

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def read_credentials(file = parent_file / "data/raw/login.txt"):
    creds = {}
    with open(file) as f:
        for line in f:
            key, value = line.strip().split("=")
            creds[key] = value
    return creds

def _set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))
    
def Hex_RGB(ip):
    return tuple(int(ip[i+1:i+3],16) for i in (0, 2, 4))

practice_track_status = {}
figures_path = '/home/kurios/Documents/f1_analysis/reports/'

prs = Presentation()

laps = practice_session.laps
stints = laps[["Driver", "Stint", "Compound", "LapNumber", "LapTime"]]
stints = stints.groupby(["Driver", "Stint", "Compound"]).agg(
    LapCount=('LapNumber', 'count'),
    FirstLap=('LapNumber', 'min'),
    LastLap=('LapNumber', 'max')
).reset_index()

air_temp_range = str((min(practice_session.weather_data['AirTemp']), max(practice_session.weather_data['AirTemp']))) + ' °C'
track_temp_range = str((min(practice_session.weather_data['TrackTemp']), max(practice_session.weather_data['TrackTemp']))) + ' °C'
humidity_range = str((min(practice_session.weather_data['Humidity']), (max(practice_session.weather_data['Humidity'])))) + ' %'

teams = fastf1.plotting.list_team_names(practice_session)

for idx, team in enumerate(teams):
    team_drivers = fastf1.plotting.get_driver_abbreviations_by_team(team, session=practice_session)

    driver_info_1 = practice_session.get_driver(team_drivers[0])
    driver_info_2 = practice_session.get_driver(team_drivers[1])
    
    #DRIVER 1 DATA
    driver_1 = driver_info_1['DriverNumber']
    driver_1_name = driver_info_1['FullName']

    epsilon = 1e-6
    practice_session.car_data[driver_1]['is_driving'] = (practice_session.car_data[driver_1]['Speed'] > epsilon) & (practice_session.car_data[driver_1]['nGear'] > 0)
    practice_session.car_data[driver_1][practice_session.car_data[driver_1]['is_driving'] == True]
    practice_session.car_data[driver_1]['block_start'] = practice_session.car_data[driver_1]['is_driving'] & (~practice_session.car_data[driver_1]['is_driving'].shift(1, fill_value=False))
    practice_session.car_data[driver_1]['block_id'] = practice_session.car_data[driver_1]['block_start'].cumsum()
    df_driving_only = practice_session.car_data[driver_1][practice_session.car_data[driver_1]['is_driving']].copy()

    if not df_driving_only.empty:
        fastest_lap = practice_session.laps.pick_drivers(driver_1).pick_fastest()
        driver_1 = driver_info_1["Abbreviation"]

        driving_block_durations = df_driving_only[df_driving_only['block_id'] > 0].groupby('block_id')['SessionTime'].apply(lambda x: x.iloc[-1] - x.iloc[0])
        min_block_duration = pd.Timedelta(minutes=1)
        filtered_driving_block_durations = driving_block_durations[driving_block_durations >= min_block_duration]
        
        total_run_time_driver_1 = str(filtered_driving_block_durations.sum())[10:-3]
        try:
            fastest_lap_driver_1 = str(fastest_lap['LapTime'])[10:-3]
        except:
            fastest_lap_driver_1 = 'No Data'

        driver_stints_driver_1 = stints[stints['Driver'] == driver_1]

        data_driver_1 = []
        try:
            for index, row in driver_stints_driver_1.iterrows():
                driver_lap = practice_session.laps.pick_drivers(driver_1).pick_laps(range(int(row['FirstLap']), int(row['LastLap'])+1)).pick_quicklaps()
                number_lap_driver_1 = int(driver_lap['LapTime'].count())
                driver_avg_lap_driver_1 = str(timedelta(seconds = driver_lap['LapTime'].dt.total_seconds().median()))[3:-3]
                compound_driver_1 = stints['Compound'][index]
                data_driver_1.append(compound_driver_1)
                data_driver_1.append(number_lap_driver_1)
                data_driver_1.append(driver_avg_lap_driver_1)
        except:
            driver_avg_lap_driver_1 = 'No Average Lap Time Found'
            compound_driver_1 = 'No Compound Found'
            number_lap_driver_1 = 'No Lap completed'
        data_driver_1_len = int(len(data_driver_1)/3)

    #DRIVER 2 DATA    
    driver_2 = driver_info_2['DriverNumber']
    driver_2_name = driver_info_2['FullName']

    epsilon = 1e-6
    practice_session.car_data[driver_2]['is_driving'] = (practice_session.car_data[driver_2]['Speed'] > epsilon) & (practice_session.car_data[driver_2]['nGear'] > 0)
    practice_session.car_data[driver_2][practice_session.car_data[driver_2]['is_driving'] == True]
    practice_session.car_data[driver_2]['block_start'] = practice_session.car_data[driver_2]['is_driving'] & (~practice_session.car_data[driver_2]['is_driving'].shift(1, fill_value=False))
    practice_session.car_data[driver_2]['block_id'] = practice_session.car_data[driver_2]['block_start'].cumsum()
    df_driving_only = practice_session.car_data[driver_2][practice_session.car_data[driver_2]['is_driving']].copy()

    if not df_driving_only.empty:
        fastest_lap = practice_session.laps.pick_drivers(driver_2).pick_fastest()
        driver_2 = driver_info_2["Abbreviation"]

        driving_block_durations = df_driving_only[df_driving_only['block_id'] > 0].groupby('block_id')['SessionTime'].apply(lambda x: x.iloc[-1] - x.iloc[0])
        min_block_duration = pd.Timedelta(minutes=1)
        filtered_driving_block_durations = driving_block_durations[driving_block_durations >= min_block_duration]
        
        total_run_time_driver_2 = str(filtered_driving_block_durations.sum())[10:-3]
        try:
            fastest_lap_driver_2 = str(fastest_lap['LapTime'])[10:-3]
        except:
            fastest_lap_driver_2 = 'No Data'

        driver_stints_driver_2 = stints[stints['Driver'] == driver_2]

        data_driver_2 = []
        try:
            for index, row in driver_stints_driver_2.iterrows():
                driver_lap = practice_session.laps.pick_drivers(driver_2).pick_laps(range(int(row['FirstLap']), int(row['LastLap'])+1)).pick_quicklaps()
                number_lap_driver_2 = int(driver_lap['LapTime'].count())
                driver_avg_lap_driver_2 = str(timedelta(seconds = driver_lap['LapTime'].dt.total_seconds().median()))[3:-3]
                compound_driver_2 = stints['Compound'][index]
                data_driver_2.append(compound_driver_2)
                data_driver_2.append(number_lap_driver_2)
                data_driver_2.append(driver_avg_lap_driver_2)
        except:
            driver_avg_lap_driver_2 = 'No Average Lap Time Found'
            compound_driver_2 = 'No Compound Found'
            number_lap_driver_2 = 'No Lap completed'
        data_driver_2_len = int(len(data_driver_2)/3)

    race_name = str(practice_session)[20:]
    
    team_logo = f'/home/kurios/Documents/f1_analysis/data/external/team_logos/{team}.png'
    
    prs.slide_width = Pt(1080)
    prs.slide_height = Pt(1350)
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)

    #BACKGROUND
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(21, 21, 30)

    #HEADER
    f1_logo = '/home/kurios/Documents/f1_analysis/data/external/team_logos/F1_75_Logo.png'
    pic = slide.shapes.add_picture(f1_logo, Pt(40), Pt(54), height=Pt(38), width= Pt(200))

    title = slide.shapes.title
    title.text = race_name
    title.top = Pt(84)
    title.left = Pt(215)
    title.width = Pt(600)

    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.name = 'Formula1 Display Bold'
        
    pic = slide.shapes.add_picture(team_logo, Pt(820), Pt(25), height= Pt(100), width=Pt(200))

    #STRUCTURE
    line1=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(150), Pt(1000), Pt(2))
    line1.line.fill.background()
    fill = line1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(244, 244, 244)

    line2=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(40), Pt(250), Pt(1000), Pt(2))
    line2.line.fill.background()
    fill = line2.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(244, 244, 244)

    line4=slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, Pt(539), Pt(250), Pt(2), Pt(1000))
    line4.line.fill.background()
    fill = line4.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(244, 244, 244)

    #REFERENCES
    txBox = slide.shapes.add_textbox(left=Pt(170), top= Pt(150), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Air Temp Range"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)
    
    txBox = slide.shapes.add_textbox(left=Pt(525), top= Pt(150), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Track Temp Range"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(910), top= Pt(150), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Humidity Range"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(70), top= Pt(250), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Driver"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)
    
    txBox = slide.shapes.add_textbox(left=Pt(1010), top= Pt(250), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Driver"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(95), top= Pt(300), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Fastest Lap"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)
    
    txBox = slide.shapes.add_textbox(left=Pt(985), top= Pt(300), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Fastest Lap"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(125), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Number Laps"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(275), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Tyre Type"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(425), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Average Lap Time"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)
    
    txBox = slide.shapes.add_textbox(left=Pt(665), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Number Laps"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(815), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Tyre Type"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    txBox = slide.shapes.add_textbox(left=Pt(965), top= Pt(400), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Average Lap Time"
    run.font.name = 'Formula1 Display Bold'
    run.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)

    if data_driver_1_len != 0:
        shape = slide.shapes.add_table(data_driver_1_len, 3,left=Pt(40), top= Pt(450), width=Pt(460), height=Pt(800))
        table = shape.table
        for i in range(0, data_driver_1_len):
            cell = table.cell(i, 0)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            y = i * 3
            run.text = f'{data_driver_1[y+1]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell = table.cell(i, 1)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f'{data_driver_1[y]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            cell = table.cell(i, 2)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f'{data_driver_1[y+2]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        shape = slide.shapes.add_table(1, 1,left=Pt(40), top= Pt(450), width=Pt(460), height=Pt(800))
        table = shape.table
        cell = table.cell(0, 0)
        fill = cell.fill
        fill.solid()
        
        fill.fore_color.rgb = RGBColor(21, 21, 30)

        text_frame = cell.text_frame
        text_frame.clear()  

        p = text_frame.add_paragraph()
        run = p.add_run()
        y = i * 3
        run.text = f'No data'
        
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if data_driver_2_len != 0:
        shape = slide.shapes.add_table(data_driver_2_len, 3,left=Pt(580), top= Pt(450), width=Pt(460), height=Pt(800))
        table = shape.table
        for i in range(0, data_driver_2_len):
            cell = table.cell(i, 0)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            y = i * 3
            run.text = f'{data_driver_2[y+1]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            cell = table.cell(i, 1)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f'{data_driver_2[y]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            cell = table.cell(i, 2)
            fill = cell.fill
            fill.solid()
            
            fill.fore_color.rgb = RGBColor(21, 21, 30)

            text_frame = cell.text_frame
            text_frame.clear()  

            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f'{data_driver_2[y+2]}'
            
            run.font.name = 'Formula1 Display Regular'
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        shape = slide.shapes.add_table(1, 1,left=Pt(580), top= Pt(450), width=Pt(460), height=Pt(800))
        table = shape.table
        cell = table.cell(0, 0)
        fill = cell.fill
        fill.solid()
        
        fill.fore_color.rgb = RGBColor(21, 21, 30)

        text_frame = cell.text_frame
        text_frame.clear()  

        p = text_frame.add_paragraph()
        run = p.add_run()
        y = i * 3
        run.text = f'No data'
        
        run.font.name = 'Formula1 Display Regular'
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    #TRACK VARIABLE
    txBox = slide.shapes.add_textbox(left=Pt(170), top=Pt(170), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = air_temp_range
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(left=Pt(525), top= Pt(170), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = track_temp_range
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(left=Pt(910), top= Pt(170), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = humidity_range
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    #DRIVERS VARIABLE
    txBox = slide.shapes.add_textbox(left=Pt(310), top=Pt(240), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = driver_1_name
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(left=Pt(770), top=Pt(240), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = driver_2_name
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(left=Pt(310), top=Pt(290), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = fastest_lap_driver_1
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(left=Pt(770), top=Pt(290), width=Pt(0), height=Pt(20))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = fastest_lap_driver_2
    run.font.name = 'Formula1 Display Regular'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

prs.save(parent_file/ report_folder / f'{race_number}_{free_practice_session}.pptx')

for file_path in report_folder.iterdir():
    if file_path.is_file() and file_path.suffix.lower() == '.pptx':
        os.chdir(report_folder)

        subprocess.run(
            [
                "/usr/bin/lowriter",
                "--headless",
                "--convert-to",
                "pdf",
                str(file_path.resolve())
            ],
            check=True
        )
        no_suffix_path = file_path.with_suffix('')
        no_suffix_path.mkdir(parents=True, exist_ok=True)
        file_path_pdf = file_path.with_suffix('.pdf')
        
        os.chdir(no_suffix_path)
        images = convert_from_path(file_path_pdf, dpi=72)
       
        for i in range(len(images)):
            images[i].save(file_path.stem + str(i) +'.jpeg', 'JPEG')
        file_path.unlink()
        file_path_pdf.unlink()

if post_option == 'Y':
    caption = ( f"Round {practice_session.event.RoundNumber} - {practice_session.event.EventName} ({practice_session.event.Country})\n"
                f"📍 {practice_session.event.Location}\n"
                f"🛠️ {practice_session} Session Recap\n" 
                f"———\n"
                f"The teams are gathering crucial data and fine-tuning their setups. Stay tuned for more updates from the {practice_session.event.OfficialEventName}!\n"
                f"#F1 #Formula1 #F1{practice_session.event.RoundNumber} #F1{practice_session.event.Country.replace(' ', '')} #F1Weekend #FreePractice #F1Practice #DataGathering #Setup")

    creds = read_credentials()
    cl = Client()
    cl.login(creds['username'], creds['password'])

    image_folder = parent_file / f"reports/{practice_session.event.RoundNumber}_{practice_session.event.EventName}_{year}/{practice_session.event.RoundNumber}_{free_practice_session}"
    image_files = []
    for report in os.listdir(image_folder):
        image_files.append(os.path.join(image_folder, report))

    cl.album_upload(
        paths=image_files,
        caption=caption
        )
