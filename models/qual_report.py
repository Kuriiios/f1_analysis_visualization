from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide=prs.slides.add_slide(title_slide_layout)
title= slide.shapes.title
title.text = 'Australian GP Qualifying Q1'
prs.save('test.pptx')